#!/usr/bin/env python3
"""Assessment 시스템 발표자료(PPTX) 생성기.
docs/operations/presentation-deck-draft.md 의 10장 구성을 네이티브 도형/표로 렌더링.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---- 팔레트 ----
NAVY   = RGBColor(0x1F, 0x2A, 0x44)
BLUE   = RGBColor(0x2E, 0x5C, 0xB5)
LBLUE  = RGBColor(0xDD, 0xE7, 0xF7)
TEAL   = RGBColor(0x1B, 0x8A, 0x8A)
ORANGE = RGBColor(0xE4, 0x7A, 0x2E)
GREY   = RGBColor(0x6B, 0x72, 0x80)
LGREY  = RGBColor(0xEE, 0xF0, 0xF3)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x22, 0x26, 0x2C)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def _set_fill(shape, color):
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color


def box(s, x, y, w, h, text="", fill=BLUE, line=None, font=14, fcolor=WHITE,
        bold=False, align=PP_ALIGN.CENTER, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        anchor=MSO_ANCHOR.MIDDLE, wrap=True):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _set_fill(sp, fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right"):
        setattr(tf, m, Inches(0.06))
    for m in ("margin_top", "margin_bottom"):
        setattr(tf, m, Inches(0.03))
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(font)
        run.font.bold = bold
        run.font.color.rgb = fcolor
        run.font.name = "Malgun Gothic"
    return sp


def textbox(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space=4):
    """runs: list of paragraphs; each paragraph = list of (text, size, color, bold)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        for (txt, size, color, bold) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Malgun Gothic"
    return tb


def arrow(s, x1, y1, x2, y2, color=GREY, width=2.0, label=None,
          dash=False, lx=None, ly=None):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    cn.shadow.inherit = False
    le = cn.line._get_or_add_ln()
    tail = le.makeelement(qn('a:tailEnd'),
                          {'type': 'triangle', 'w': 'med', 'h': 'med'})
    le.append(tail)
    if dash:
        d = le.makeelement(qn('a:prstDash'), {'val': 'dash'})
        le.append(d)
    if label:
        lx = lx if lx is not None else (x1 + x2) / 2 - 0.7
        ly = ly if ly is not None else (y1 + y2) / 2 - 0.18
        t = textbox(s, lx, ly, 1.6, 0.3,
                    [[(label, 10, color, True)]], align=PP_ALIGN.CENTER)
        t.text_frame.paragraphs[0].runs[0].font.name = "Malgun Gothic"
    return cn


def header(s, num, title, sub=None):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.05))
    _set_fill(bar, NAVY)
    bar.line.fill.background()
    bar.shadow.inherit = False
    box(s, 0.35, 0.16, 0.72, 0.72, str(num), fill=BLUE, font=26, bold=True,
        shape=MSO_SHAPE.OVAL)
    runs = [[(title, 26, WHITE, True)]]
    if sub:
        runs.append([(sub, 13, RGBColor(0xC7, 0xD2, 0xE8), False)])
    textbox(s, 1.25, 0.12, 11.6, 0.85, runs, anchor=MSO_ANCHOR.MIDDLE)


def footer(s, n):
    textbox(s, 11.7, 7.08, 1.5, 0.3,
            [[("Assessment · %02d" % n, 9, GREY, False)]], align=PP_ALIGN.RIGHT)


# =====================================================================
# 0. 표지
# =====================================================================
s = slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
_set_fill(bg, NAVY); bg.line.fill.background(); bg.shadow.inherit = False
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.55), SW, Inches(0.12))
_set_fill(accent, ORANGE); accent.line.fill.background(); accent.shadow.inherit = False
textbox(s, 1.0, 2.3, 11.3, 2.2, [
    [("Assessment 평가·진단 시스템", 40, WHITE, True)],
    [("테스트 환경 구성 & 배포 아키텍처", 24, RGBColor(0xAD,0xC4,0xE8), False)],
])
textbox(s, 1.0, 4.85, 11.3, 1.5, [
    [("기술팀 내부 공유", 16, RGBColor(0xC7,0xD2,0xE8), False)],
    [("발표자: taewon  ·  2026-06-18", 14, GREY, False)],
    [("상세 커버리지 매트릭스는 별첨 문서 참조", 12, GREY, False)],
])

# =====================================================================
# 1. 시스템 구조 한눈에
# =====================================================================
s = slide()
header(s, 1, "한눈에 보기 — 시스템 구조", "무엇을 만들었는가: 기능 코드가 아니라 '배포 인프라'")
# 외부
box(s, 0.5, 2.4, 1.9, 1.0, "사내망\n사용자", fill=GREY, font=14, bold=True)
box(s, 0.5, 4.6, 1.9, 1.0, "Bastion\n배포·운영 거점", fill=DARK, font=13, bold=True)
# 폐쇄망 컨테이너
cont = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.1), Inches(1.55),
                          Inches(9.6), Inches(5.2))
_set_fill(cont, LGREY); cont.line.color.rgb = GREY; cont.line.width = Pt(1.25)
cont.shadow.inherit = False
textbox(s, 3.3, 1.62, 5, 0.4, [[("OpenStack 폐쇄망", 13, GREY, True)]])
box(s, 3.6, 2.4, 4.0, 1.9,
    "Engine VM\n\napi · consumer\npostgres · rabbitmq · redis\n= docker compose 단일 스택",
    fill=BLUE, font=14, bold=True)
box(s, 8.3, 2.4, 4.0, 0.95, "AI VM\nOllama + diagnostic-worker",
    fill=TEAL, font=13, bold=True)
box(s, 8.3, 3.55, 4.0, 0.75, "—  AMQP / Ollama :11434  —", fill=None,
    line=TEAL, fcolor=TEAL, font=11)
box(s, 3.6, 4.7, 8.7, 1.7,
    "Agent 플릿  ·  39대\nLinux 31  +  Windows 8\n(점검 대상 서버 시뮬레이션)",
    fill=ORANGE, font=15, bold=True)
arrow(s, 2.4, 2.9, 3.6, 3.1, color=BLUE, width=2.5, label="HTTP :8000", ly=2.55)
arrow(s, 7.6, 5.3, 5.6, 4.3, color=ORANGE, width=2.5, label="AMQP :5672", lx=5.7, ly=4.75)
arrow(s, 8.3, 3.0, 7.6, 3.0, color=TEAL, width=2.0)
arrow(s, 1.45, 4.6, 1.45, 3.4, color=GREY, width=1.5)
arrow(s, 2.4, 5.1, 3.1, 5.1, color=GREY, width=1.5, dash=True, label="SSH 배포", ly=5.25)
textbox(s, 0.5, 6.55, 12.3, 0.7, [[
    ("VM 3종 — ① 평가 엔진(1) · ② AI 진단(1) · ③ agent 플릿(39, 점검 대상 시뮬레이션)",
     13, DARK, True)]])
footer(s, 1)


# =====================================================================
# 2. 사용자 시나리오 (고객 관점 여정 — 시퀀스 대신 4단계 흐름)
# =====================================================================
s = slide()
header(s, 2, "고객 사용 시나리오",
       "고객 관점 — 에이전트 하나로 서버 인식부터 원격 작업·AI 진단까지")
steps = [
    ("1", "에이전트 설치", "고객 서버에\nassessment-agent 배포", ORANGE),
    ("2", "서버 자동 인식", "agent가 스스로 보고 →\n대시보드에 서버 등장", BLUE),
    ("3", "원격 작업 발행", "ZDM install 등\n작업을 원격으로 실행", TEAL),
    ("4", "AI 진단 확인", "작업 결과에 LLM\n진단 코멘트가 붙음", NAVY),
]
x0, y, bw, bh, gap = 0.55, 2.4, 2.85, 2.5, 0.42
for i, (num, title, desc, c) in enumerate(steps):
    x = x0 + i * (bw + gap)
    box(s, x, y, bw, bh, "", fill=WHITE, line=c)
    box(s, x+bw/2-0.45, y-0.45, 0.9, 0.9, num, fill=c, font=28, bold=True,
        shape=MSO_SHAPE.OVAL)
    textbox(s, x+0.15, y+0.7, bw-0.3, bh-0.8, [
        [(title, 18, c, True)],
        [("", 6, DARK, False)],
        [(desc, 13, DARK, False)],
    ], align=PP_ALIGN.CENTER)
    if i < 3:
        ax = x + bw + 0.05
        arrow(s, ax, y+bh/2, ax+gap-0.1, y+bh/2, color=GREY, width=2.5)
box(s, 0.55, 5.4, 12.25, 1.4,
    "고객 입장에선 '에이전트 설치 → 대시보드에 서버가 뜬다'가 전부.\n"
    "단순 모니터링을 넘어 원격 작업 발행 + AI 진단까지가 이 제품의 차별점.",
    fill=LGREY, line=GREY, fcolor=DARK, font=15, bold=True, align=PP_ALIGN.LEFT)
footer(s, 2)

# =====================================================================
# 3. 테스트 환경 — 구성 개요 (무엇을·왜)
# =====================================================================
s = slide()
header(s, 3, "테스트 환경 — 구성 개요",
       "고객 서버는 제각각 → '실제로 띄워서' 동작을 실측하는 환경")
# 좌: 목표
box(s, 0.6, 1.55, 5.7, 0.6, "무엇을 만들었나", fill=BLUE, font=16, bold=True)
textbox(s, 0.8, 2.35, 5.5, 4.6, [
    [("● 39대 멀티-OS 테스트 플릿", 16, DARK, True)],
    [("   Linux 31대 + Windows 8대", 13, GREY, False)],
    [("", 7, DARK, False)],
    [("● 각 VM에 로컬 서비스 7종 설치", 16, DARK, True)],
    [("   db·cache·mq·web·app·container·monitor", 13, GREY, False)],
    [("   = 실제 고객 워크로드 재현", 13, GREY, False)],
    [("", 7, DARK, False)],
    [("● noise 레이어로 장애·부하 시뮬레이션", 16, DARK, True)],
    [("", 7, DARK, False)],
    [("● 서브넷 8개 분산 + dual-homed", 16, DARK, True)],
    [("   멀티 NIC·서브넷 간 통신까지 검증", 13, GREY, False)],
])
# 우: 왜
box(s, 6.7, 1.55, 6.1, 0.6, "왜 이렇게까지", fill=ORANGE, font=16, bold=True)
textbox(s, 6.9, 2.35, 5.8, 4.6, [
    [("고객 서버 OS가 제각각이다.", 16, DARK, True)],
    [("", 6, DARK, False)],
    [("빌드 매트릭스로 '될 것이다' 추론이 아니라,", 14, DARK, False)],
    [("각 OS 이미지에 인스턴스를 실제로 띄우고", 14, DARK, True)],
    [("agent를 배포해 직접 확인한다.", 14, DARK, True)],
    [("", 8, DARK, False)],
    [("검증 핵심 2가지:", 15, BLUE, True)],
    [("  ✓ MQ 발행 — 엔진이 서버를 인식하는가", 14, DARK, False)],
    [("  ✓ ZDM install — 원격 작업이 수행되는가", 14, DARK, False)],
    [("", 8, DARK, False)],
    [("⇒ 현실 조건(워크로드+노이즈)에서", 14, GREEN, True)],
    [("   agent가 버티는지까지 본다", 14, GREEN, True)],
])
footer(s, 3)

# =====================================================================
# 4. 테스트 환경 — OS 매트릭스 (왜 39대 / 검증 3축)
# =====================================================================
s = slide()
header(s, 4, "테스트 환경 — OS 커버리지",
       "검증 축이 3개: OS 종류 × firmware × Python 버전")
rows = [
    ("계열", "커버 범위", "대수"),
    ("Linux", "Debian·Ubuntu·Rocky·Alma\nCentOS·Oracle·Amazon", "31"),
    ("Windows", "2003 · 2008 · 2012 ~ 2025", "8"),
]
heights = [0.5, 1.0, 0.7]
gtbl = s.shapes.add_table(3, 3, Inches(0.5), Inches(1.5), Inches(6.5),
                          Inches(sum(heights))).table
gtbl.columns[0].width = Inches(1.4)
gtbl.columns[1].width = Inches(4.0)
gtbl.columns[2].width = Inches(1.1)
for r in range(3):
    gtbl.rows[r].height = Inches(heights[r])
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        c = gtbl.cell(ri, ci)
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY if ri == 0 else (LGREY if ri % 2 else WHITE)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.08); c.margin_top = Inches(0.02)
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci == 2 else (PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT)
        for j, ln in enumerate(val.split("\n")):
            pp = p if j == 0 else c.text_frame.add_paragraph()
            pp.alignment = p.alignment
            run = pp.add_run(); run.text = ln
            run.font.size = Pt(13 if ri == 0 else 13)
            run.font.bold = (ri == 0 or ci == 0)
            run.font.color.rgb = WHITE if ri == 0 else DARK
            run.font.name = "Malgun Gothic"
box(s, 0.5, 3.75, 6.5, 1.15,
    "목적: 각 이미지에 인스턴스를 띄우고 agent 배포 →\nMQ 발행 + ZDM install 이 실제로 되는지 실측",
    fill=NAVY, font=14, bold=True, align=PP_ALIGN.LEFT)
box(s, 0.5, 5.15, 6.5, 1.5,
    "왜 단순 빌드 매트릭스가 아니라 실측인가\n\n"
    "같은 바이너리도 OS·firmware·Python 버전에 따라\n부팅·ZDM 동작이 달라지기 때문",
    fill=LGREY, line=GREY, fcolor=DARK, font=13, bold=False, align=PP_ALIGN.LEFT)
# 우측 3축
box(s, 7.3, 1.5, 5.5, 0.5, "검증 축 3개", fill=BLUE, font=16, bold=True)
axes = [
    ("OS 종류", "Linux 7개 배포판 + Windows 8세대"),
    ("Firmware", "같은 OS도 BIOS / UEFI 따로 검증\n(부팅·ZDM 동작 상이)"),
    ("Python 버전", "py<3.7 OS는 정적 바이너리\n수동 배포 경로로 분기"),
]
yy = 2.2
for t, d in axes:
    box(s, 7.3, yy, 5.5, 1.4, "", fill=WHITE, line=ORANGE)
    textbox(s, 7.5, yy+0.18, 5.1, 1.1, [
        [("● " + t, 16, ORANGE, True)],
        [(d, 13, DARK, False)],
    ])
    yy += 1.55
footer(s, 4)

# =====================================================================
# 5. 테스트 환경 — 서브넷 토폴로지
# =====================================================================
s = slide()
header(s, 5, "테스트 환경 — 서브넷 토폴로지",
       "primary는 고정 / secondary는 분산 + dual-homed")
box(s, 0.5, 1.35, 12.3, 0.7,
    "primary NIC — 전 agent 'agent-subnet' 고정  (engine MQ · bastion 연결, 필수)",
    fill=BLUE, font=14, bold=True)
labels = ["test-net\n4대", "-02\n4대", "-03\n4대", "-04\n4대",
          "-05\n4대", "-06\n4대", "-07\n4대", "-08\n3대"]
n = len(labels)
x0, y0, bw, bh, gap = 0.55, 3.1, 1.35, 1.0, 0.18
xs = []
for i, lb in enumerate(labels):
    x = x0 + i * (bw + gap)
    xs.append(x)
    box(s, x, y0, bw, bh, lb, fill=ORANGE if i == 0 else TEAL, font=12, bold=True)
for i in range(n - 1):
    arrow(s, xs[i]+bw, y0+0.5, xs[i+1], y0+0.5, color=GREY, width=1.5)
wrap_y = y0 + bh + 0.55
arrow(s, xs[-1]+bw/2, y0+bh, xs[-1]+bw/2, wrap_y, color=ORANGE, width=1.5)
ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                            Inches(xs[-1]+bw/2), Inches(wrap_y),
                            Inches(xs[0]+bw/2), Inches(wrap_y))
ln.line.color.rgb = ORANGE; ln.line.width = Pt(1.5); ln.shadow.inherit = False
arrow(s, xs[0]+bw/2, wrap_y, xs[0]+bw/2, y0+bh, color=ORANGE, width=1.5)
textbox(s, 5.0, wrap_y-0.02, 3.5, 0.3,
        [[("wrap-around (dual-homed)", 11, ORANGE, True)]], align=PP_ALIGN.CENTER)
textbox(s, 0.6, 5.5, 12.2, 1.6, [
    [("● secondary NIC — Linux 31대를 4대씩 8개 서브넷에 분산", 14, DARK, True)],
    [("● dual-homed — 각 그룹 첫 대는 인접 서브넷에도 연결", 14, DARK, True)],
    [("   → 멀티 NIC · 서브넷 간 통신 시나리오까지 검증", 13, GREY, False)],
    [("   신규 서브넷 6개는 internal-only 스택에서 Terraform 생성 (ADR-0013)", 12, GREY, False)],
])
footer(s, 5)

# =====================================================================
# 6. 테스트 환경 — 워크로드 시뮬레이션
# =====================================================================
s = slide()
header(s, 6, "테스트 환경 — 워크로드 시뮬레이션",
       "빈 VM이 아니라 '실제 고객 서버처럼'")
cont = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5),
                          Inches(7.5), Inches(5.1))
_set_fill(cont, LGREY); cont.line.color.rgb = GREY; cont.line.width = Pt(1.5)
cont.shadow.inherit = False
textbox(s, 0.9, 1.6, 5, 0.4, [[("agent VM 1대", 14, GREY, True)]])
box(s, 1.1, 2.2, 6.7, 0.8, "assessment-agent", fill=BLUE, font=15, bold=True)
box(s, 1.1, 3.25, 6.7, 0.45, "로컬 서비스 7종 = 점검 대상 워크로드",
    fill=None, fcolor=TEAL, font=13, bold=True, align=PP_ALIGN.LEFT)
svcs = ["db", "cache", "mq", "web", "app", "container", "monitor"]
sx, sw2 = 1.1, 0.92
for i, sv in enumerate(svcs):
    box(s, sx + i*(sw2+0.02), 3.75, sw2, 0.7, sv, fill=TEAL, font=12, bold=True)
box(s, 1.1, 4.75, 6.7, 1.4,
    "noise 레이어\n재시작 · offline · 부하 시뮬레이션",
    fill=ORANGE, font=14, bold=True)
arrow(s, 4.45, 4.75, 4.45, 4.45, color=ORANGE, width=2.0)
textbox(s, 8.5, 1.9, 4.4, 4.5, [
    [("왜 이렇게?", 18, BLUE, True)],
    [("", 6, DARK, False)],
    [("실제 고객 서버에는 DB·캐시·웹·컨테이너가", 13, DARK, False)],
    [("다 돌고 있다.", 13, DARK, False)],
    [("", 6, DARK, False)],
    [("→ 7종 로컬 서비스를 모두 설치해", 13, DARK, True)],
    [("   현실 워크로드를 재현", 13, DARK, False)],
    [("", 6, DARK, False)],
    [("→ noise role로 장애·부하까지 흉내", 13, DARK, True)],
    [("", 6, DARK, False)],
    [("⇒ agent가 현실 조건에서 버티는지 검증", 13, GREEN, True)],
])
footer(s, 6)

# =====================================================================
# 7. 테스트 환경 — 어떻게 구축했나 (폐쇄망 + Terraform→Ansible)
# =====================================================================
s = slide()
header(s, 7, "테스트 환경 — 어떻게 구축했나",
       "폐쇄망 제약 위에서 Terraform → Ansible 파이프라인으로 자동 구성")
# 파이프라인
box(s, 0.7, 1.55, 3.6, 1.0, "Horizon\nnet · router · keypair\n(수동 1회)",
    fill=GREY, font=13, bold=True)
box(s, 4.85, 1.55, 3.6, 1.0, "Terraform\nSG · VM · port · volume",
    fill=BLUE, font=13, bold=True)
box(s, 9.0, 1.55, 3.6, 1.0, "Ansible\n바이너리 배포 · 서비스 · noise",
    fill=TEAL, font=13, bold=True)
arrow(s, 4.3, 2.05, 4.85, 2.05, color=DARK, width=2.5)
arrow(s, 8.45, 2.05, 9.0, 2.05, color=DARK, width=2.5)
# 폐쇄망
box(s, 0.7, 3.0, 12.0, 0.55, "폐쇄망 제약 — VM은 외부 인터넷 직접 접근 불가",
    fill=ORANGE, font=15, bold=True)
box(s, 0.9, 3.85, 2.7, 1.1, "인터넷\nReleases·GHCR", fill=GREY, font=13, bold=True)
box(s, 5.3, 3.85, 2.7, 1.1, "Bastion", fill=DARK, font=15, bold=True)
box(s, 9.7, 3.85, 2.7, 1.1, "Agent VM", fill=BLUE, font=14, bold=True)
arrow(s, 3.6, 4.4, 5.3, 4.4, color=GREEN, width=2.0, label="✓ 대신 받음", ly=4.08)
arrow(s, 8.0, 4.4, 9.7, 4.4, color=BLUE, width=2.5, label="바이너리 주입", ly=4.08)
textbox(s, 0.7, 5.25, 12.1, 1.7, [
    [("● 자원(VM·SG·volume)은 Terraform 선언적 IaC / VM 내부 설정·secret은 Ansible",
      14, DARK, True)],
    [("● bastion이 agent 바이너리를 대신 받아 files/에 사전 복사 → Ansible로 주입",
      14, DARK, True)],
    [("● python<3.7 OS는 정적 바이너리 수동 배포 / Windows는 WinRM + win_copy (ADR-0007)",
      13, GREY, False)],
])
footer(s, 7)

# =====================================================================
# 8. 검증 기준 & 현황
# =====================================================================
s = slide()
header(s, 8, "검증 기준 & 현황", "상세 결과는 별첨 커버리지 매트릭스 문서")
rows = [
    ("검증 항목", "합격 기준"),
    ("MQ 발행", "agent 기동 후 engine server_inventory에 신규 인식"),
    ("ZDM install", "install 작업 발행 → agent 수행 → status success"),
    ("권한 모델", "생성 유저 / sudo / systemd User= 기록"),
]
heights = [0.5, 0.7, 0.7, 0.7]
tbl = s.shapes.add_table(4, 2, Inches(0.5), Inches(1.45), Inches(8.0),
                         Inches(sum(heights))).table
tbl.columns[0].width = Inches(2.3)
tbl.columns[1].width = Inches(5.7)
for r in range(4):
    tbl.rows[r].height = Inches(heights[r])
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        c = tbl.cell(ri, ci)
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY if ri == 0 else (LGREY if ri % 2 else WHITE)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.1)
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = val
        run.font.size = Pt(14 if ri == 0 else 13)
        run.font.bold = (ri == 0 or ci == 0)
        run.font.color.rgb = WHITE if ri == 0 else DARK
        run.font.name = "Malgun Gothic"
box(s, 8.8, 1.45, 4.05, 0.5, "알려진 예외 (by-design)", fill=ORANGE,
    font=14, bold=True)
textbox(s, 8.9, 2.05, 4.0, 3.0, [
    [("● Windows ZDM install", 13, DARK, True)],
    [("  엔진 패키지가 Linux 전용 → 구조적 불가", 12, GREY, False)],
    [("  (별도 지원 작업 필요)", 12, GREY, False)],
    [("", 6, DARK, False)],
    [("● RHEL 계열", 13, DARK, True)],
    [("  SELinux enforcing이 install.sh 차단", 12, GREY, False)],
    [("  → permissive 적용 (ADR-0012)", 12, GREY, False)],
])
box(s, 0.5, 4.6, 8.0, 2.0,
    "📎  OS별 BIOS/UEFI 전수 결과는\n별첨 커버리지 매트릭스 문서 참조\n(agent-binary-image-coverage.md)",
    fill=LGREY, line=GREY, fcolor=DARK, font=14, bold=True, align=PP_ALIGN.LEFT)
footer(s, 8)

# =====================================================================
# 9. 마무리
# =====================================================================
s = slide()
header(s, 9, "마무리 — 현황 & 남은 과제", None)
box(s, 0.6, 1.5, 5.9, 0.6, "✅  완료", fill=GREEN, font=18, bold=True)
textbox(s, 0.8, 2.3, 5.7, 3.5, [
    [("OS 매트릭스·서브넷 토폴로지 확정", 14, DARK, True)],
    [("(Linux 31 + Windows 8)", 12, GREY, False)],
    [("", 6, DARK, False)],
    [("Terraform 모델·tfvars 작성, validate 통과", 14, DARK, True)],
    [("", 6, DARK, False)],
    [("engine v0.7.0 가동 중", 14, DARK, True)],
])
box(s, 6.9, 1.5, 5.9, 0.6, "☐  남은 작업", fill=ORANGE, font=18, bold=True)
textbox(s, 7.1, 2.3, 5.7, 3.5, [
    [("network 스택 apply → 서브넷 6개 생성", 14, DARK, True)],
    [("", 6, DARK, False)],
    [("agent 스택 apply → 39대 부팅", 14, DARK, True)],
    [("", 6, DARK, False)],
    [("전수 배포 → MQ·ZDM 검증 → 매트릭스 기입", 14, DARK, True)],
])
box(s, 0.6, 6.0, 12.2, 0.9, "감사합니다   ·   Q & A", fill=NAVY, font=22, bold=True)
footer(s, 9)

out = "/home/debian/assessment-infra/docs/operations/assessment-presentation.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
